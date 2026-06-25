#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
객체지향프로그래밍 중간고사 - 문항별 답안 취합 스크립트
=======================================================
총 4문항 기준으로 모든 학생의 답안을 문항별로 모아
'문항별_답안_취합.pdf' 파일을 생성합니다.

[필요 패키지 설치]
  pip install pymupdf pillow python-pptx python-docx reportlab pdfplumber

[LibreOffice 설치 필요 (docx/hwp 변환용)]
  https://www.libreoffice.org/download/download/

[실행 방법]
  python 문항별_취합_실행.py
"""

import os
import re
import sys
import shutil
import subprocess
import zipfile
import argparse
from pathlib import Path

# ─────────────────────────────────────────────
# 커맨드라인 인자 파싱
# ─────────────────────────────────────────────
def parse_args():
    parser = argparse.ArgumentParser(
        description="객체지향프로그래밍 중간고사 문항별 답안 취합",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
사용 예시:
  python 문항별_취합_실행.py
  python 문항별_취합_실행.py C:\\Users\\prof\\중간고사
  python 문항별_취합_실행.py "C:\\제출 폴더" --output "C:\\결과\\취합.pdf" --questions 5
        """,
    )
    parser.add_argument(
        "input_dir",
        nargs="?",                          # 생략 가능
        default=None,
        help="답안 파일이 있는 폴더 경로 (기본값: 스크립트와 같은 폴더)",
    )
    parser.add_argument(
        "--output", "-o",
        default=None,
        help="출력 폴더 경로 (기본값: input_dir/취합결과/). PDF 파일명은 자동 생성됨.",
    )
    parser.add_argument(
        "--questions", "-q",
        type=int,
        default=4,
        help="총 문항 수 (기본값: 4)",
    )
    return parser.parse_args()

# ─────────────────────────────────────────────
# 설정  ← parse_args() 결과로 채워짐 (main에서)
# ─────────────────────────────────────────────
SCRIPT_DIR    = Path(__file__).parent
INPUT_DIR     = None   # main()에서 설정
OUTPUT_PDF    = None   # main()에서 설정
TMP_DIR       = None   # main()에서 설정
NUM_QUESTIONS = 4      # main()에서 설정

# LibreOffice 경로 (설치 위치에 따라 수정)
LIBREOFFICE_PATHS = [
    r"C:\Program Files\LibreOffice\program\soffice.exe",
    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    "soffice",   # PATH에 등록된 경우
]

# 문항 번호 감지 패턴 — 명확한 헤더 형식만 허용
# "[1]", "(1):", "1." 등 답안 본문의 소항목과 혼동되는 패턴은 제외
QUESTION_PATTERNS = [
    (r'[★☆■□▶▷◆◇]\s*([1-4])\s*번', 1),   # ★1번, ■2번
    (r'문[제항]\s*([1-4])\s*[번.]?', 1),    # 문제1, 문항2번
    (r'^\s*([1-4])\s*번\s*[.:\s]', 1),      # 줄 시작 "1번." "2번:"
    (r'^[Qq][.\s]*([1-4])\b', 1),            # 줄 시작 Q1, Q.2
]

# ─────────────────────────────────────────────
# 유틸리티
# ─────────────────────────────────────────────

def find_libreoffice():
    for path in LIBREOFFICE_PATHS:
        if os.path.isfile(path) or shutil.which(path):
            return path
    return None

def extract_student_name(filename: str) -> str:
    """파일명 맨 앞 3글자를 학생 이름으로 반환."""
    stem = Path(filename).stem
    return stem[:3]

def detect_question_from_text(text: str) -> list:
    """
    텍스트 앞부분(첫 10줄)에서만 문항 번호를 감지.
    답안 본문 속 번호(소제목, 참조 등)를 문항 번호로 오인하지 않도록
    페이지 상단 헤더 영역만 검사한다.
    """
    header = '\n'.join(text.splitlines()[:10])
    found = set()
    for pattern, group in QUESTION_PATTERNS:
        for m in re.finditer(pattern, header, re.MULTILINE | re.IGNORECASE):
            q = int(m.group(group))
            if 1 <= q <= NUM_QUESTIONS:
                found.add(q)
    return sorted(found)


def split_text_by_questions(text: str) -> dict:
    """
    전체 텍스트에서 문항 헤더를 찾아 문항별 구간으로 분리.
    반환: {q: text_segment} — 2개 이상 감지된 경우에만 유효한 값 반환.
    학생이 4문항을 한 페이지에 연속으로 작성한 경우에 사용.
    """
    first_pos = {}
    for pattern, group in QUESTION_PATTERNS:
        for m in re.finditer(pattern, text, re.MULTILINE | re.IGNORECASE):
            q = int(m.group(group))
            if 1 <= q <= NUM_QUESTIONS:
                if q not in first_pos or m.start() < first_pos[q]:
                    first_pos[q] = m.start()

    if len(first_pos) < 2:
        return {}

    boundaries = sorted(first_pos.items(), key=lambda x: x[1])
    result = {}
    for idx, (q, pos) in enumerate(boundaries):
        end = boundaries[idx + 1][1] if idx + 1 < len(boundaries) else len(text)
        result[q] = text[pos:end].strip()
    return result


def assign_questions_by_position(num_pages: int) -> dict:
    """
    텍스트 감지 실패 시: 페이지 위치 기반으로 문항 배정.
    - 페이지 수 == 문항 수: 페이지 N → 문항 N
    - 페이지 수 > 문항 수: 균등 분할
    - 페이지 수 < 문항 수: 모든 페이지를 1번 문항으로
    """
    assignment = {}
    if num_pages == 0:
        return assignment
    if num_pages == NUM_QUESTIONS:
        for i in range(num_pages):
            assignment[i] = [i + 1]
    elif num_pages > NUM_QUESTIONS:
        # 앞쪽 페이지들을 문항에 배정
        pages_per_q = num_pages // NUM_QUESTIONS
        for i in range(num_pages):
            q = min(i // pages_per_q + 1, NUM_QUESTIONS)
            assignment[i] = [q]
    else:
        # 페이지 수가 부족 → 각 페이지가 순서대로 문항에 해당
        for i in range(num_pages):
            assignment[i] = [i + 1]
    return assignment


# ─────────────────────────────────────────────
# 파일 변환: → 이미지 목록
# ─────────────────────────────────────────────

def _open_fitz_silently(pdf_path: Path):
    """
    MuPDF 경고 메시지를 stderr에 출력하지 않고 PDF를 엽니다.
    'No common ancestor in structure tree' 같은 구조 경고는
    Python 예외가 아니라 C 레벨 stderr 출력이므로
    fitz.TOOLS 설정 + stderr 리다이렉션으로 억제합니다.
    """
    import fitz, os, sys, io

    # PyMuPDF 내장 경고 출력 비활성화
    try:
        fitz.TOOLS.mupdf_display_errors(False)
    except AttributeError:
        pass  # 구버전 호환

    # C 레벨 stderr도 억제 (os-level redirect)
    devnull_fd = os.open(os.devnull, os.O_WRONLY)
    old_stderr_fd = os.dup(2)
    os.dup2(devnull_fd, 2)
    os.close(devnull_fd)

    doc = None
    try:
        doc = fitz.open(str(pdf_path))
    except Exception:
        pass
    finally:
        # stderr 복원
        os.dup2(old_stderr_fd, 2)
        os.close(old_stderr_fd)

    return doc


def pdf_to_images(pdf_path: Path, out_dir: Path, prefix: str) -> list:
    """PDF를 페이지별 PNG 이미지로 변환. 이미지 경로 목록 반환."""
    import fitz  # PyMuPDF

    # ── 1차 시도: MuPDF (경고 억제) ─────────────────
    doc = _open_fitz_silently(pdf_path)

    # ── 2차 시도: stream 모드 ────────────────────────
    if doc is None:
        try:
            raw = Path(pdf_path).read_bytes()
            doc = fitz.open(stream=raw, filetype="pdf")
        except Exception as e:
            print(f"  ⚠ fitz stream 열기 실패: {e}")

    # ── fitz로 렌더링 ────────────────────────────────
    if doc is not None:
        images = []
        mat = fitz.Matrix(2.0, 2.0)
        for i in range(len(doc)):
            try:
                page = doc.load_page(i)
                pix  = page.get_pixmap(matrix=mat, alpha=False)
                img_path = out_dir / f"{prefix}_page{i+1:03d}.png"
                pix.save(str(img_path))
                images.append((img_path, page.get_text()))
            except Exception as e:
                print(f"  ⚠ {i+1}페이지 렌더링 실패 (건너뜀): {e}")
        doc.close()
        if images:
            return images
        print(f"  ⚠ fitz 렌더링 결과 없음 → pypdf로 재시도")

    # ── 3차 시도: pypdf + Pillow 렌더링 ─────────────
    try:
        from pypdf import PdfReader
        from PIL import Image as PILImage
        import io

        reader = PdfReader(str(pdf_path), strict=False)
        images = []
        for i, page in enumerate(reader.pages):
            # 텍스트 추출
            try:
                text = page.extract_text() or ""
            except Exception:
                text = ""

            # 페이지를 이미지로 렌더링 (pypdf는 렌더링 불가 → 텍스트 이미지 fallback)
            img_path = text_to_image(text, out_dir, f"{prefix}_page{i+1:03d}")
            images.append((img_path, text))

        if images:
            print(f"  ℹ pypdf fallback: {len(images)}페이지 텍스트 추출")
            return images
    except Exception as e:
        print(f"  ⚠ pypdf 실패: {e}")

    # ── 4차 시도: pdfplumber ─────────────────────────
    try:
        import pdfplumber
        images = []
        with pdfplumber.open(str(pdf_path)) as plumb:
            for i, page in enumerate(plumb.pages):
                text = page.extract_text() or ""
                img_path = text_to_image(text, out_dir, f"{prefix}_page{i+1:03d}")
                images.append((img_path, text))
        if images:
            print(f"  ℹ pdfplumber fallback: {len(images)}페이지 텍스트 추출")
            return images
    except Exception as e:
        print(f"  ✗ pdfplumber도 실패: {e}")

    return []


def pptx_to_images(pptx_path: Path, out_dir: Path, prefix: str) -> list:
    """PPTX 슬라이드를 이미지로 변환 (LibreOffice 경유)."""
    lo = find_libreoffice()
    if lo:
        # LibreOffice로 PDF 변환 후 이미지화
        tmp_pdf = out_dir / f"{prefix}_converted.pdf"
        try:
            subprocess.run(
                [lo, "--headless", "--convert-to", "pdf",
                 "--outdir", str(out_dir), str(pptx_path)],
                timeout=60, capture_output=True
            )
            # 변환된 PDF 찾기
            converted = list(out_dir.glob(f"{pptx_path.stem}*.pdf"))
            if converted:
                converted[0].rename(tmp_pdf)
                imgs = pdf_to_images(tmp_pdf, out_dir, prefix)
                tmp_pdf.unlink(missing_ok=True)
                return imgs
        except Exception as e:
            print(f"  ⚠ LibreOffice 변환 실패 ({prefix}): {e}")

    # Fallback: python-pptx로 슬라이드 텍스트만 추출 → 텍스트를 이미지로 렌더
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        results = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
            slide_text = "\n".join(texts)
            # 텍스트를 흰 배경 이미지에 렌더링
            img_path = text_to_image(slide_text, out_dir, f"{prefix}_slide{i+1:03d}")
            results.append((img_path, slide_text))
        return results
    except Exception as e:
        print(f"  ✗ PPTX 처리 실패 ({prefix}): {e}")
        return []


def docx_to_images(docx_path: Path, out_dir: Path, prefix: str) -> list:
    """DOCX를 이미지로 변환 (LibreOffice 경유)."""
    lo = find_libreoffice()
    if lo:
        tmp_pdf = out_dir / f"{prefix}_converted.pdf"
        try:
            subprocess.run(
                [lo, "--headless", "--convert-to", "pdf",
                 "--outdir", str(out_dir), str(docx_path)],
                timeout=60, capture_output=True
            )
            converted = list(out_dir.glob(f"{docx_path.stem}*.pdf"))
            if converted:
                converted[0].rename(tmp_pdf)
                imgs = pdf_to_images(tmp_pdf, out_dir, prefix)
                tmp_pdf.unlink(missing_ok=True)
                return imgs
        except Exception as e:
            print(f"  ⚠ LibreOffice 변환 실패 ({prefix}): {e}")

    # Fallback: python-docx로 텍스트 추출 (단락 + 표 모두 포함)
    try:
        from docx import Document
        from docx.oxml.ns import qn

        doc = Document(str(docx_path))

        # 문서 본문 순서대로 단락과 표를 함께 순회
        body = doc.element.body
        lines = []
        for child in body:
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag == 'p':
                text = ''.join(
                    node.text for node in child.iter(qn('w:t'))
                    if node.text
                )
                if text.strip():
                    lines.append(text)
            elif tag == 'tbl':
                for tr in child.iter(qn('w:tr')):
                    seen_cell = set()
                    cell_texts = []
                    for tc in tr.iter(qn('w:tc')):
                        t = ''.join(
                            node.text for node in tc.iter(qn('w:t'))
                            if node.text
                        ).strip()
                        if t and t not in seen_cell:
                            seen_cell.add(t)
                            cell_texts.append(t)
                    if cell_texts:
                        lines.append(' | '.join(cell_texts))

        full_text = '\n'.join(lines)
        img_path = text_to_image(full_text, out_dir, f"{prefix}_page001")
        return [(img_path, full_text)]
    except Exception as e:
        print(f"  ✗ DOCX 처리 실패 ({prefix}): {e}")
        return []


def _hwpx_parse_text(hwpx_path: Path) -> list:
    """
    HWPX 파일(ZIP+XML 구조)에서 섹션별 텍스트를 추출.
    HWPX는 추가 라이브러리 없이 zipfile + ElementTree로 파싱 가능.
    반환: [section0_text, section1_text, ...]
    """
    import zipfile
    import xml.etree.ElementTree as ET

    sections = []
    try:
        with zipfile.ZipFile(str(hwpx_path)) as z:
            # Contents/section0.xml, section1.xml … 순서대로 읽기
            section_files = sorted(
                f for f in z.namelist()
                if re.match(r'Contents/section\d+\.xml', f)
            )
            if not section_files:
                # 일부 HWPX는 경로가 다를 수 있음
                section_files = sorted(
                    f for f in z.namelist()
                    if 'section' in f.lower() and f.endswith('.xml')
                )
            for sf in section_files:
                with z.open(sf) as f:
                    content = f.read().decode('utf-8', errors='ignore')
                root = ET.fromstring(content)
                # HWP XML에서 텍스트 노드 <hp:t> 추출
                # 네임스페이스 무관하게 tag 뒷부분이 't'인 것 수집
                texts = []
                for elem in root.iter():
                    local = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                    if local == 't' and elem.text and elem.text.strip():
                        texts.append(elem.text)
                if texts:
                    sections.append('\n'.join(texts))
    except Exception as e:
        print(f"  ⚠ HWPX XML 파싱 오류: {e}")
    return sections


def _hwp_via_hancom_com(hwp_path: Path, out_dir: Path, prefix: str) -> list:
    """
    Hancom Office(한컴오피스)가 설치된 Windows에서 COM 자동화로 PDF 변환.
    pywin32 패키지 필요: pip install pywin32
    """
    try:
        import win32com.client
        tmp_pdf = out_dir / f"{prefix}_hancom.pdf"
        hwp = None
        try:
            hwp = win32com.client.Dispatch("HWPFrame.HwpObject")
            # 보안 모듈 등록 (파일 경로 접근 허용)
            try:
                hwp.RegisterModule("FilePathCheckDLL", "FilePathCheckerModule")
            except Exception:
                pass
            hwp.Open(str(hwp_path.resolve()), "HWP", "forceopen:true")
            hwp.SaveAs(str(tmp_pdf.resolve()), "PDF")
            hwp.Quit()
            hwp = None
        except Exception as e:
            if hwp:
                try: hwp.Quit()
                except: pass
            raise e

        if tmp_pdf.exists() and tmp_pdf.stat().st_size > 0:
            imgs = pdf_to_images(tmp_pdf, out_dir, prefix)
            tmp_pdf.unlink(missing_ok=True)
            return imgs
    except ImportError:
        pass  # pywin32 미설치
    except Exception as e:
        print(f"  ⚠ Hancom COM 변환 실패: {e}")
    return []


def _hwp_via_pyhwp(hwp_path: Path, out_dir: Path, prefix: str) -> list:
    """
    pyhwp 라이브러리로 HWP 텍스트 추출.
    설치: pip install pyhwp
    """
    try:
        import hwp5  # noqa — 존재 확인용
    except ImportError:
        return []

    try:
        txt_path = out_dir / f"{prefix}_pyhwp.txt"
        result = subprocess.run(
            [sys.executable, '-m', 'hwp5.hwp5txt',
             '--output', str(txt_path), str(hwp_path)],
            capture_output=True, timeout=30
        )
        if txt_path.exists() and txt_path.stat().st_size > 0:
            text = txt_path.read_text(encoding='utf-8', errors='ignore')
            txt_path.unlink(missing_ok=True)
            # 텍스트를 이미지로 변환 (페이지 구분 없음 → 통짜 1장)
            img_path = text_to_image(text, out_dir, f"{prefix}_page001")
            print(f"  ℹ pyhwp 텍스트 추출 성공")
            return [(img_path, text)]
    except Exception as e:
        print(f"  ⚠ pyhwp 변환 실패: {e}")
    return []


def _hwp_via_libreoffice(hwp_path: Path, out_dir: Path, prefix: str) -> list:
    """LibreOffice로 HWP/HWPX → PDF 변환."""
    lo = find_libreoffice()
    if not lo:
        return []
    try:
        subprocess.run(
            [lo, "--headless", "--convert-to", "pdf",
             "--outdir", str(out_dir), str(hwp_path)],
            timeout=60, capture_output=True
        )
        converted = list(out_dir.glob(f"{hwp_path.stem}*.pdf"))
        if converted:
            tmp_pdf = out_dir / f"{prefix}_lo.pdf"
            converted[0].rename(tmp_pdf)
            imgs = pdf_to_images(tmp_pdf, out_dir, prefix)
            tmp_pdf.unlink(missing_ok=True)
            if imgs:
                return imgs
    except Exception as e:
        print(f"  ⚠ LibreOffice 변환 실패: {e}")
    return []


def _parse_hwp_records(data: bytes) -> str:
    """HWP 5.0 레코드 스트림(압축 해제 후)에서 텍스트 추출."""
    import struct
    HWPTAG_PARA_TEXT = 67
    texts = []
    pos = 0
    while pos + 4 <= len(data):
        header = struct.unpack_from('<I', data, pos)[0]
        tag_id = header & 0x3FF
        size   = (header >> 20) & 0xFFF
        pos += 4
        if size == 0xFFF:
            if pos + 4 > len(data):
                break
            size = struct.unpack_from('<I', data, pos)[0]
            pos += 4
        if pos + size > len(data):
            break
        if tag_id == HWPTAG_PARA_TEXT:
            chunk = data[pos:pos + size]
            chars = []
            for ci in range(0, len(chunk) - 1, 2):
                code = struct.unpack_from('<H', chunk, ci)[0]
                if code == 13:
                    chars.append('\n')
                elif code >= 32:
                    chars.append(chr(code))
            texts.append(''.join(chars))
        pos += size
    return '\n'.join(texts)


def _hwp_via_olefile(hwp_path: Path, out_dir: Path, prefix: str) -> list:
    """
    olefile로 HWP 5.0 바이너리에서 텍스트 직접 추출.
    설치: pip install olefile
    """
    try:
        import olefile
        import zlib
    except ImportError:
        return []

    try:
        if not olefile.isOleFile(str(hwp_path)):
            return []

        ole = olefile.OleFileIO(str(hwp_path))
        sections = []
        i = 0
        while ole.exists(f'BodyText/Section{i}'):
            raw = ole.openstream(f'BodyText/Section{i}').read()
            # HWP 5.0 섹션 스트림: raw deflate 압축
            decompressed = None
            for wbits in (-15, 15, 47):
                try:
                    decompressed = zlib.decompress(raw, wbits)
                    break
                except Exception:
                    continue
            if decompressed is None:
                decompressed = raw
            text = _parse_hwp_records(decompressed)
            if text.strip():
                sections.append(text)
            i += 1
        ole.close()

        if sections:
            imgs = []
            for j, text in enumerate(sections):
                img_path = text_to_image(text, out_dir, f"{prefix}_sec{j+1:03d}")
                imgs.append((img_path, text))
            print(f"  ✓ olefile HWP 파싱: {len(imgs)}섹션")
            return imgs
    except Exception as e:
        print(f"  ⚠ olefile HWP 파싱 실패: {e}")
    return []


def hwp_to_images(hwp_path: Path, out_dir: Path, prefix: str) -> list:
    """HWP / HWPX 파일을 이미지로 변환 (다단계 시도)."""
    ext = hwp_path.suffix.lower()

    # ── HWPX: ZIP+XML 직접 파싱 ──────────────────────
    if ext == '.hwpx':
        sections = _hwpx_parse_text(hwp_path)
        if sections:
            imgs = []
            for i, text in enumerate(sections):
                img_path = text_to_image(text, out_dir, f"{prefix}_sec{i+1:03d}")
                imgs.append((img_path, text))
            print(f"  ✓ HWPX 직접 파싱: {len(imgs)}섹션")
            return imgs

    # ── 1차: Hancom Office COM 자동화 ────────────────
    print(f"  → Hancom COM 시도...")
    imgs = _hwp_via_hancom_com(hwp_path, out_dir, prefix)
    if imgs:
        print(f"  ✓ Hancom COM 변환 성공: {len(imgs)}페이지")
        return imgs

    # ── 2차: pyhwp 텍스트 추출 ───────────────────────
    print(f"  → pyhwp 시도...")
    imgs = _hwp_via_pyhwp(hwp_path, out_dir, prefix)
    if imgs:
        return imgs

    # ── 3차: LibreOffice ──────────────────────────────
    print(f"  → LibreOffice 시도...")
    imgs = _hwp_via_libreoffice(hwp_path, out_dir, prefix)
    if imgs:
        print(f"  ✓ LibreOffice 변환 성공: {len(imgs)}페이지")
        return imgs

    # ── 4차: olefile 직접 파싱 ────────────────────────
    print(f"  → olefile 직접 파싱 시도...")
    imgs = _hwp_via_olefile(hwp_path, out_dir, prefix)
    if imgs:
        return imgs

    print(f"  ✗ HWP 처리 실패 (모든 방법 소진): {hwp_path.name}")
    print(f"    해결 방법: pip install olefile  또는  LibreOffice 설치")
    return []


def text_to_image(text: str, out_dir: Path, prefix: str) -> Path:
    """텍스트를 흰 배경 PNG 이미지로 변환 (LibreOffice 없을 때 fallback)."""
    from PIL import Image, ImageDraw, ImageFont

    # Windows 한글 폰트 전체 경로 순서대로 시도
    FONT_CANDIDATES = [
        (r"C:\Windows\Fonts\malgun.ttf",   28),
        (r"C:\Windows\Fonts\malgunbd.ttf", 28),
        (r"C:\Windows\Fonts\gulim.ttc",    26),
        (r"C:\Windows\Fonts\batang.ttc",   26),
    ]

    font = None
    for font_path, size in FONT_CANDIDATES:
        if os.path.isfile(font_path):
            try:
                font = ImageFont.truetype(font_path, size)
                break
            except Exception:
                continue
    if font is None:
        font = ImageFont.load_default()

    img = Image.new("RGB", (1240, 1754), color="white")  # A4 @ 150dpi
    draw = ImageDraw.Draw(img)

    # 긴 텍스트는 줄바꿈 처리
    lines = []
    for raw_line in text.splitlines():
        while len(raw_line) > 60:
            lines.append(raw_line[:60])
            raw_line = raw_line[60:]
        lines.append(raw_line)
        if len(lines) > 80:  # 최대 80줄
            lines.append("... (이하 생략)")
            break

    y = 60
    for line in lines:
        draw.text((60, y), line, fill="black", font=font)
        y += 34
        if y > 1700:
            break

    img_path = out_dir / f"{prefix}.png"
    img.save(str(img_path))
    return img_path


# ─────────────────────────────────────────────
# 문항별 이미지 분류
# ─────────────────────────────────────────────

def _primary_question(text: str):
    """
    페이지 첫 15줄에서 문항 번호를 하나만 반환 (가장 먼저 등장한 것).
    15줄을 쓰는 이유: 학생 정보(이름·학번·제출일 등) 헤더가 5줄 이상을
    차지하는 경우가 많아 5줄만 보면 문항 번호를 놓침.
    first-match 방식이므로 본문에서 다른 문항 번호를 참조하더라도
    헤더의 문항 번호가 먼저 등장하면 올바르게 처리된다.
    """
    header = '\n'.join(text.splitlines()[:15])
    best_pos, best_q = None, None
    for pattern, group in QUESTION_PATTERNS:
        for m in re.finditer(pattern, header, re.MULTILINE | re.IGNORECASE):
            q = int(m.group(group))
            if 1 <= q <= NUM_QUESTIONS:
                if best_pos is None or m.start() < best_pos:
                    best_pos, best_q = m.start(), q
    return best_q


def classify_pages(pages: list, filename: str) -> dict:
    """
    pages = [(img_path, text), ...]
    반환: {question_num: [img_path, ...]}

    핵심 원칙: 한 페이지는 반드시 하나의 문항에만 배정.
    텍스트에서 가장 먼저 등장하는 문항 번호(primary)를 사용하므로,
    답안 본문에 다른 문항 번호가 언급돼도 중복 배정되지 않는다.
    """
    result = {q: [] for q in range(1, NUM_QUESTIONS + 1)}
    n = len(pages)
    if n == 0:
        return result

    # 각 페이지의 주 문항 번호 (정확히 하나, 또는 None)
    page_primary = {}
    for i, (img_path, text) in enumerate(pages):
        q = _primary_question(text)
        if q is not None:
            page_primary[i] = q

    detected = set(page_primary.values())

    if len(detected) == NUM_QUESTIONS:
        # 모든 문항이 감지된 경우에만 텍스트 기반 사용.
        # 일부만 감지된 경우({1,4} 등)는 감지 안 된 페이지가
        # carry-forward로 1번에 쏠리므로 위치 기반으로 폴백.
        last_q = page_primary[min(page_primary.keys())]
        for i, (img_path, _) in enumerate(pages):
            if i in page_primary:
                last_q = page_primary[i]
            result[last_q].append(img_path)
        print(f"  ℹ 텍스트 기반 분류 {sorted(detected)}: {filename}")
    else:
        # 위치 기반 배정
        reason = "전체 문항 미감지" if detected else "문항 감지 실패"
        print(f"  ℹ 위치 기반 배정 ({n}페이지, {reason}): {filename}")
        assignment = assign_questions_by_position(n)
        for i, (img_path, _) in enumerate(pages):
            for q in assignment.get(i, [1]):
                result[q].append(img_path)

    return result


# ─────────────────────────────────────────────
# DOCX 표 기반 문항별 답안 추출
# ─────────────────────────────────────────────

def _q_from_header_cell(text: str):
    """
    표 첫 행 텍스트에서 문항 번호를 추출.
    QUESTION_PATTERNS보다 유연하게: '1번', '문항1', '1', 'Q1', '[1]' 등 처리.
    """
    text = text.strip()
    if not text:
        return None

    q = _primary_question(text)
    if q is not None:
        return q

    HEADER_PATTERNS = [
        r'([1-9])\s*번',
        r'문[제항]\s*([1-9])',
        r'[Qq]\s*([1-9])',
        r'\[([1-9])\]',
        r'\(([1-9])\)',
        r'^\s*([1-9])\s*$',
    ]
    for pat in HEADER_PATTERNS:
        m = re.search(pat, text, re.IGNORECASE)
        if m:
            n = int(m.group(1))
            if 1 <= n <= NUM_QUESTIONS:
                return n
    return None




def _compose_answer_image(header: str, items: list,
                           out_dir: Path, filename: str) -> Path:
    """
    헤더 + items(str 또는 PIL.Image 혼합)를 문서 순서대로 세로 합성해 PNG 저장.
    str → 텍스트 줄,  PIL.Image → 이미지 블록.
    """
    from PIL import Image as PILImage, ImageDraw, ImageFont

    WIDTH     = 1240
    MARGIN    = 16           # 좌우 여백 최소화
    LINE_H = 46           # 폰트 34 + 줄간격 12
    img_w  = WIDTH - 2 * MARGIN

    FONT_CANDIDATES = [
        (r"C:\Windows\Fonts\malgun.ttf",   34),
        (r"C:\Windows\Fonts\malgunbd.ttf", 34),
        (r"C:\Windows\Fonts\gulim.ttc",    32),
        (r"C:\Windows\Fonts\batang.ttc",   32),
    ]
    font = None
    for fp, sz in FONT_CANDIDATES:
        if os.path.isfile(fp):
            try:
                font = ImageFont.truetype(fp, sz)
                break
            except Exception:
                continue
    if font is None:
        font = ImageFont.load_default()

    def _text_width(text: str) -> int:
        try:
            return int(font.getlength(text))
        except AttributeError:
            return font.getsize(text)[0]

    def _wrap(text: str) -> list:
        """실제 픽셀 폭 기준으로 줄바꿈. 탭은 공백 4개로 대체."""
        text = text.replace('\t', '    ')
        if not text:
            return ['']
        max_w = img_w
        lines = []
        current = ''
        for ch in text:
            probe = current + ch
            if _text_width(probe) > max_w and current:
                lines.append(current)
                current = ch
            else:
                current = probe
        lines.append(current)
        return lines

    # items를 렌더 블록으로 변환: ('text', [줄...]) 또는 ('image', PIL.Image)
    render = []
    if header:
        render.append(('text', [header]))

    text_count = 0
    for item in items:
        if isinstance(item, str):
            if text_count >= 120:
                if not any(k == 'text' and '(이하 생략)' in v[-1]
                           for k, v in render[-1:]):
                    render.append(('text', ['... (이하 생략)']))
                continue
            wrapped = _wrap(item)
            text_count += len(wrapped)
            render.append(('text', wrapped))
        else:
            iw, ih = item.size
            new_h = max(1, int(ih * img_w / iw))
            render.append(('image', item.resize((img_w, new_h), PILImage.LANCZOS)))

    # 전체 높이 계산
    total_h = MARGIN
    for kind, content in render:
        if kind == 'text':
            total_h += len(content) * LINE_H
        else:
            total_h += content.height + MARGIN
    total_h += MARGIN

    canvas = PILImage.new('RGB', (WIDTH, max(1, total_h)), 'white')
    draw   = ImageDraw.Draw(canvas)

    y = MARGIN
    for kind, content in render:
        if kind == 'text':
            for line in content:
                draw.text((MARGIN, y), line, fill='black', font=font)
                y += LINE_H
        else:
            canvas.paste(content, (MARGIN, y))
            y += content.height + MARGIN

    out_path = out_dir / filename
    canvas.save(str(out_path))
    return out_path


def extract_tables_from_docx(docx_path: Path, out_dir: Path, prefix: str) -> dict:
    """
    DOCX 표에서 문항별 답안(텍스트 + 이미지)을 추출.
    ZIP을 직접 열어 rels, document.xml, 이미지를 모두 ET로 처리한다.
    python-docx(lxml)와 ET의 네임스페이스 표현 차이를 회피하기 위해
    lxml을 사용하지 않고 ET 단독으로 파싱한다.
    """
    import xml.etree.ElementTree as ET
    import io as _io
    from PIL import Image as PILImage

    def _local(tag: str) -> str:
        return tag.split('}')[-1] if '}' in tag else tag

    rId_to_img: dict = {}
    doc_root = None

    # ── ZIP에서 이미지 + document.xml 동시 로드 ──────
    try:
        with zipfile.ZipFile(str(docx_path)) as z:
            znames = set(z.namelist())

            # 이미지 관계 로드
            rels_path = 'word/_rels/document.xml.rels'
            if rels_path in znames:
                for rel in ET.fromstring(z.read(rels_path)):
                    if not rel.get('Type', '').endswith('/image'):
                        continue
                    rid    = rel.get('Id', '').strip()
                    target = rel.get('Target', '').strip()
                    zp     = (f"word/{target}" if not target.startswith('/')
                              else target.lstrip('/'))
                    if zp not in znames:
                        continue
                    try:
                        raw = z.read(zp)
                        img = PILImage.open(_io.BytesIO(raw))
                        img.load()
                        rId_to_img[rid] = img.convert('RGB')
                        print(f"  ✓ 이미지: {rid} ← {zp} {img.size}")
                    except Exception as e:
                        print(f"  ⚠ 이미지 실패 ({zp}): {e}")

            # document.xml 로드
            doc_xml_path = 'word/document.xml'
            if doc_xml_path in znames:
                doc_root = ET.fromstring(z.read(doc_xml_path))

    except Exception as e:
        print(f"  ⚠ ZIP 실패: {e}")
        return {}

    print(f"  ℹ 이미지 {len(rId_to_img)}개: {list(rId_to_img.keys())}")

    if doc_root is None:
        print(f"  ⚠ document.xml 없음")
        return {}

    # ── document.xml에서 표 직접 탐색 ────────────────
    # 최상위 tbl 요소만 수집 (중첩 표 제외)
    body = next((el for el in doc_root.iter() if _local(el.tag) == 'body'), doc_root)
    tables = [el for el in body if _local(el.tag) == 'tbl']
    print(f"  ℹ 표 {len(tables)}개")

    if not tables:
        return {}

    result = {}
    for ti, tbl in enumerate(tables):
        rows = [el for el in tbl if _local(el.tag) == 'tr']
        if not rows:
            continue

        # 첫 행 텍스트 (중복 없이)
        first_texts = []
        seen_ft = set()
        for tc in rows[0].iter():
            if _local(tc.tag) == 't' and tc.text:
                t = tc.text.strip()
                if t and t not in seen_ft:
                    seen_ft.add(t)
                    first_texts.append(t)
        first_text = ' '.join(first_texts)
        print(f"  [표 {ti+1}] '{first_text[:50]}'", end='  ')

        q = _q_from_header_cell(first_text)
        if q is None:
            print("→ 번호 없음")
            continue
        print(f"→ {q}번")

        answer_items: list = []   # str(텍스트 줄) 또는 PIL.Image 혼합, 문서 순서
        seen_rids:    set  = set()

        def _iter_row_cells(row_elem):
            """행 내 w:tc를 재귀 탐색. w:sdt 래퍼도 통과. 중첩 w:tbl은 건너뜀."""
            for child in row_elem:
                loc = _local(child.tag)
                if loc == 'tc':
                    yield child
                elif loc != 'tbl':
                    yield from _iter_row_cells(child)

        def _iter_cell_paras(elem):
            """셀 내 w:p를 재귀 탐색. 중첩 w:tbl은 건너뜀."""
            for child in elem:
                loc = _local(child.tag)
                if loc == 'p':
                    yield child
                elif loc != 'tbl':
                    yield from _iter_cell_paras(child)

        def _para_items(para):
            """단락 내 텍스트 줄·이미지를 문서 순서대로 반환."""
            items: list = []
            buf:   list = []
            for el in para.iter():
                loc = _local(el.tag)
                if loc == 't' and el.text is not None:
                    buf.append(el.text)
                elif loc == 'tab':
                    buf.append('\t')
                elif loc == 'br':
                    items.append(''.join(buf))
                    buf = []
                elif loc == 'blip':
                    # 이미지 앞 텍스트 먼저 flush
                    line = ''.join(buf)
                    if line:
                        items.append(line)
                    buf = []
                    rid = None
                    for ak, av in el.attrib.items():
                        if _local(ak) == 'embed':
                            rid = av.strip()
                            break
                    if rid and rid not in seen_rids and rid in rId_to_img:
                        seen_rids.add(rid)
                        items.append(rId_to_img[rid])
            items.append(''.join(buf))   # 단락 끝 나머지 텍스트
            return items

        for row in rows[1:]:
            seen_tc_ids: set = set()
            for tc in _iter_row_cells(row):
                tc_id = id(tc)
                if tc_id in seen_tc_ids:
                    continue
                seen_tc_ids.add(tc_id)
                for para in _iter_cell_paras(tc):
                    answer_items.extend(_para_items(para))

        has_text  = any(isinstance(x, str) and x.strip() for x in answer_items)
        has_image = any(not isinstance(x, str) for x in answer_items)
        if not has_text and not has_image:
            print(f"    → 내용 없음")
            continue

        n_txt = sum(1 for x in answer_items if isinstance(x, str))
        n_img = sum(1 for x in answer_items if not isinstance(x, str))
        print(f"    → 텍스트 {n_txt}줄, 이미지 {n_img}개")
        out_img = _compose_answer_image(
            f"[{q}번 문항]", answer_items,
            out_dir, f"{prefix}_q{q}.png"
        )
        result[q] = out_img

    return result


# ─────────────────────────────────────────────
# 메인 처리
# ─────────────────────────────────────────────

def _rmtree_onerror(func, path, _excinfo):
    """rmtree 오류 핸들러: 읽기 전용 속성 제거 후 재시도 (Windows WinError 5 대응)."""
    import stat
    try:
        os.chmod(path, stat.S_IWRITE)
        func(path)
    except Exception:
        pass


def process_all_submissions():
    """모든 제출 파일을 처리하여 {question: [(name, [img_paths])]} 반환"""

    # 임시 폴더 초기화
    if TMP_DIR.exists():
        shutil.rmtree(TMP_DIR, onerror=_rmtree_onerror)
    TMP_DIR.mkdir(parents=True)

    # ZIP 압축 해제
    zip_files = list(INPUT_DIR.glob("*.zip"))
    for zf in zip_files:
        print(f"📦 ZIP 압축 해제: {zf.name}")
        with zipfile.ZipFile(zf) as z:
            z.extractall(TMP_DIR / "zip_extracted")

    # 처리 대상 파일 수집
    # - hwp+pdf 중복 제출자는 pdf 우선
    all_files = list(INPUT_DIR.glob("*.*"))
    all_files += list(TMP_DIR.glob("zip_extracted/**/*.*"))

    # 이 스크립트가 생성하는 출력 파일명 집합 (재처리 방지)
    out_dir = OUTPUT_PDF.parent
    output_filenames = {OUTPUT_PDF.name}

    # 출력 디렉토리가 입력 디렉토리 내부에 있을 때만 해당 경로로 필터링.
    # out_dir이 INPUT_DIR의 부모인 경우 relative_to가 모든 파일에 성공해
    # 전체 파일이 제외되는 버그를 방지.
    try:
        out_dir.relative_to(INPUT_DIR)
        out_dir_is_inside_input = True
    except ValueError:
        out_dir_is_inside_input = False

    # 학생별로 파일 그루핑 (파일명 앞부분 기준)
    student_files = {}
    SUPPORTED = {".pdf", ".pptx", ".docx", ".hwp", ".hwpx"}

    for f in all_files:
        if f.suffix.lower() not in SUPPORTED:
            continue
        # 출력 디렉토리가 입력 디렉토리 안에 있을 때만 해당 폴더 내 파일 제외
        if out_dir_is_inside_input:
            try:
                f.relative_to(out_dir)
                continue
            except ValueError:
                pass
        # 출력 파일 자체와 이름이 일치하면 제외
        if f.name in output_filenames:
            continue

        name = extract_student_name(f.name)
        if name not in student_files:
            student_files[name] = []
        student_files[name].append(f)

    # 중복 제출: PDF 우선, hwp 후순위
    priority_files = {}
    for name, files in student_files.items():
        exts = {f.suffix.lower() for f in files}
        if ".pdf" in exts:
            # PDF 있으면 PDF만 사용
            selected = [f for f in files if f.suffix.lower() == ".pdf"]
        elif ".docx" in exts:
            selected = [f for f in files if f.suffix.lower() == ".docx"]
        elif ".pptx" in exts:
            selected = [f for f in files if f.suffix.lower() == ".pptx"]
        else:
            selected = files
        priority_files[name] = selected[0]  # 파일 1개만 선택

    print(f"\n📋 처리 대상 학생 수: {len(priority_files)}명")

    # 문항별 결과: {q: [(student_name, img_path), ...]}
    question_data = {q: [] for q in range(1, NUM_QUESTIONS + 1)}
    failed = []

    for student_name, f in sorted(priority_files.items()):
        print(f"\n👤 처리 중: {student_name} ({f.name})")
        ext = f.suffix.lower()
        student_tmp = TMP_DIR / student_name
        student_tmp.mkdir(exist_ok=True)
        prefix = student_name

        try:
            if ext == ".pdf":
                pages = pdf_to_images(f, student_tmp, prefix)
            elif ext == ".pptx":
                pages = pptx_to_images(f, student_tmp, prefix)
            elif ext == ".docx":
                table_imgs = extract_tables_from_docx(f, student_tmp, prefix)
                if table_imgs:
                    print(f"  ✓ DOCX 표 기반 추출: {sorted(table_imgs.keys())}번 문항")
                    for q, img in table_imgs.items():
                        question_data[q].append((student_name, img))
                    continue
                pages = docx_to_images(f, student_tmp, prefix)
            elif ext in (".hwp", ".hwpx"):
                pages = hwp_to_images(f, student_tmp, prefix)
            else:
                pages = []

            if not pages:
                print(f"  ✗ 변환 실패: {f.name}")
                failed.append(f.name)
                continue

            print(f"  ✓ 페이지 수: {len(pages)}")

            # 페이지 수가 문항 수보다 적으면 전체 텍스트에서 문항 경계를 찾아 분리
            if len(pages) < NUM_QUESTIONS:
                combined_text = '\n'.join(t for _, t in pages)
                splits = split_text_by_questions(combined_text)
                if len(splits) >= 2:
                    print(f"  ✓ 텍스트 문항 분리 성공: {sorted(splits.keys())}번")
                    for q, q_text in splits.items():
                        img = text_to_image(q_text, student_tmp, f"{prefix}_q{q}")
                        question_data[q].append((student_name, img))
                    continue  # classify_pages 건너뜀
                else:
                    print(f"  ℹ 텍스트 분리 불가 → 위치 기반 배정")

            classified = classify_pages(pages, f.name)

            for q, img_paths in classified.items():
                for img in img_paths:
                    question_data[q].append((student_name, img))

        except Exception as e:
            print(f"  ✗ 오류 ({f.name}): {e}")
            failed.append(f.name)

    return question_data, failed


# ─────────────────────────────────────────────
# PDF 생성
# ─────────────────────────────────────────────

def register_korean_font():
    """
    Windows 시스템 폰트(맑은 고딕)를 reportlab에 등록.
    반환: 등록된 폰트 이름 (없으면 기본 Helvetica 사용)
    """
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    candidates = [
        (r"C:\Windows\Fonts\malgun.ttf",   "MalgunGothic"),
        (r"C:\Windows\Fonts\malgunbd.ttf", "MalgunGothicBold"),
        (r"C:\Windows\Fonts\gulim.ttc",    "Gulim"),
        (r"C:\Windows\Fonts\batang.ttc",   "Batang"),
    ]

    registered = None
    for font_path, font_name in candidates:
        if os.path.isfile(font_path):
            try:
                # .ttc 파일은 subfontIndex=0 지정
                if font_path.endswith(".ttc"):
                    pdfmetrics.registerFont(TTFont(font_name, font_path, subfontIndex=0))
                else:
                    pdfmetrics.registerFont(TTFont(font_name, font_path))
                print(f"  ✅ 한글 폰트 등록: {font_name} ({font_path})")
                if registered is None:
                    registered = font_name
            except Exception as e:
                print(f"  ⚠ 폰트 등록 실패 ({font_name}): {e}")

    if registered is None:
        print("  ⚠ 한글 폰트를 찾지 못했습니다. 한글 텍스트가 깨질 수 있습니다.")
        return "Helvetica"
    return registered


def create_output_pdf(question_data: dict, failed: list) -> Path:
    """모든 문항 답안을 하나의 PDF 파일로 생성."""
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import (
        SimpleDocTemplate, Image as RLImage,
        Paragraph, Spacer, PageBreak, HRFlowable
    )
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER
    from PIL import Image as PILImage

    KO_FONT = register_korean_font()

    def ko_style(name, **kwargs):
        return ParagraphStyle(name, fontName=KO_FONT, **kwargs)

    style_title = ko_style(
        "KoTitle", fontSize=26, alignment=TA_CENTER,
        spaceAfter=12, leading=34,
    )
    style_sub = ko_style(
        "KoSub", fontSize=16, alignment=TA_CENTER,
        textColor=colors.HexColor("#1a4f8a"), spaceAfter=10,
    )
    style_count = ko_style(
        "KoCount", fontSize=13, alignment=TA_CENTER,
        textColor=colors.HexColor("#555555"), spaceAfter=6,
    )
    style_warn = ko_style(
        "KoWarn", fontSize=10, alignment=TA_CENTER,
        textColor=colors.red, spaceAfter=4,
    )
    style_q_title = ko_style(
        "KoQTitle", fontSize=20, alignment=TA_CENTER,
        textColor=colors.HexColor("#1a4f8a"),
        spaceAfter=10, spaceBefore=6, leading=28,
    )
    style_student = ko_style(
        "KoStudent", fontSize=13,
        textColor=colors.HexColor("#1a3a6e"),
        backColor=colors.HexColor("#dce8f8"),
        borderPadding=(8, 5, 8, 5),
        spaceAfter=10, leading=18,
    )
    style_missing = ko_style(
        "KoMissing", fontSize=12,
        textColor=colors.HexColor("#888888"),
    )
    style_err = ko_style(
        "KoErr", fontSize=10, textColor=colors.red,
    )

    W, H = A4
    content_width = W - 1.4*cm   # 좌우 각 0.7cm
    out_path = OUTPUT_PDF
    out_path.parent.mkdir(parents=True, exist_ok=True)

    print(f"\n📄 PDF 생성 중: {out_path.name}")

    doc = SimpleDocTemplate(
        str(out_path),
        pagesize=A4,
        leftMargin=0.7*cm, rightMargin=0.7*cm,
        topMargin=1.2*cm, bottomMargin=1.2*cm,
    )

    story = []

    # ── 표지 ─────────────────────────────────────
    total_entries = sum(len(v) for v in question_data.values())
    story.append(Spacer(1, 4*cm))
    story.append(Paragraph("객체지향프로그래밍 중간고사", style_title))
    story.append(Spacer(1, 0.3*cm))
    story.append(Paragraph("문항별 답안 취합본", style_sub))
    story.append(Spacer(1, 0.8*cm))
    for q in range(1, NUM_QUESTIONS + 1):
        cnt = len(question_data.get(q, []))
        story.append(Paragraph(f"{q}번 문항: {cnt}명", style_count))
    story.append(Spacer(1, 0.3*cm))
    story.append(Paragraph(f"전체 답안: {total_entries}건", style_count))

    if failed:
        failed_str = ", ".join(Path(f).name for f in failed)
        story.append(Spacer(1, 0.3*cm))
        story.append(Paragraph(
            f"변환 실패 ({len(failed)}개): {failed_str}", style_warn
        ))
    story.append(PageBreak())

    # ── 문항별 답안 ───────────────────────────────
    for q in range(1, NUM_QUESTIONS + 1):
        entries = question_data.get(q, [])
        print(f"  {q}번 문항: {len(entries)}명")

        story.append(Paragraph(
            f"{q}번 문항  —  답안 취합  ({len(entries)}명)",
            style_q_title
        ))
        story.append(HRFlowable(
            width="100%", thickness=2,
            color=colors.HexColor("#1a4f8a"),
            spaceAfter=12,
        ))

        if not entries:
            story.append(Paragraph("이 문항에 해당하는 답안이 없습니다.", style_missing))
            story.append(PageBreak())
            continue

        for student_name, img_path in entries:
            story.append(Paragraph(f"[ {student_name} ]", style_student))
            try:
                pil_img = PILImage.open(str(img_path))
                iw, ih = pil_img.size
                ratio = ih / iw
                draw_w = content_width
                draw_h = draw_w * ratio
                max_h = H - 3*cm
                if draw_h > max_h:
                    draw_h = max_h
                    draw_w = draw_h / ratio
                story.append(RLImage(str(img_path), width=draw_w, height=draw_h))
            except Exception as e:
                story.append(Paragraph(f"[이미지 로드 실패: {e}]", style_err))

            story.append(HRFlowable(
                width="100%", thickness=1,
                color=colors.HexColor("#bbbbbb"),
                spaceAfter=10,
            ))
            story.append(PageBreak())

    doc.build(story)
    print(f"  ✅ 저장: {out_path.name}")
    return out_path


# ─────────────────────────────────────────────
# 실행
# ─────────────────────────────────────────────

if __name__ == "__main__":
    args = parse_args()

    # ── 전역 설정값 확정 ──────────────────────────

    INPUT_DIR     = Path(args.input_dir) if args.input_dir else SCRIPT_DIR
    NUM_QUESTIONS = args.questions
    # --output은 폴더 경로. 미지정 시 입력 폴더 내 취합결과/ 하위 폴더 사용.
    out_dir       = Path(args.output) if args.output else INPUT_DIR / "취합결과"
    OUTPUT_PDF    = out_dir / "문항별_답안_취합.pdf"
    TMP_DIR       = INPUT_DIR / "_tmp_images"

    # ── 유효성 검사 ───────────────────────────────
    if not INPUT_DIR.exists():
        print(f"❌ 폴더를 찾을 수 없습니다: {INPUT_DIR}")
        sys.exit(1)

    print("=" * 60)
    print("  객체지향프로그래밍 중간고사 문항별 취합 시작")
    print("=" * 60)
    print(f"  입력 폴더 : {INPUT_DIR}")
    print(f"  출력 폴더 : {out_dir}")
    print(f"  출력 파일 : {OUTPUT_PDF.name}")
    print(f"  총 문항 수: {NUM_QUESTIONS}문항")
    print("=" * 60)

    # ── 패키지 확인 ───────────────────────────────
    missing = []
    for pkg, imp in [("pymupdf", "fitz"), ("Pillow", "PIL"),
                     ("python-pptx", "pptx"), ("python-docx", "docx"),
                     ("reportlab", "reportlab"), ("pdfplumber", "pdfplumber"),
                     ("pypdf", "pypdf"), ("olefile", "olefile")]:
        try:
            __import__(imp)
        except ImportError:
            missing.append(pkg)

    if missing:
        print(f"\n❌ 누락된 패키지: {', '.join(missing)}")
        print("   다음 명령어로 설치하세요:")
        print(f"   pip install {' '.join(missing)}")
        sys.exit(1)

    lo = find_libreoffice()
    if lo:
        print(f"✅ LibreOffice 발견: {lo}")
    else:
        print("⚠  LibreOffice 미설치 — docx/hwp 변환에 제한이 있습니다.")
        print("   https://www.libreoffice.org/download/ 에서 설치 권장")

    print()
    question_data, failed = process_all_submissions()
    out_path = create_output_pdf(question_data, failed)

    print("\n" + "=" * 60)
    print("  처리 완료!")
    for q, entries in question_data.items():
        print(f"  {q}번 문항: {len(entries)}개 답안")
    print(f"\n  생성된 PDF: {out_path.name}")
    if failed:
        print(f"\n  변환 실패 파일 ({len(failed)}개):")
        for f in failed:
            print(f"    - {f}")
    print("=" * 60)
